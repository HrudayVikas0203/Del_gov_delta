from pathlib import Path
from textwrap import shorten

from openpyxl import Workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from sqlalchemy.orm import Session

from app.core.config import get_settings
from app.models.delivery import Account, Project
from app.models.people import Employee
from app.models.status import GeneratedReport, ReportFormat, WeeklyStatus
from app.schemas.common import LLMSelection


BRAND_BLUE = "#1D4ED8"
BRAND_NAVY = "#0F172A"
INK_SOFT = "#475569"
SUCCESS = "#16A34A"
WARNING = "#D97706"
DANGER = "#DC2626"
LIGHT_BLUE = "#EFF6FF"
BORDER_BLUE = "#BFDBFE"


def _report_path(report: GeneratedReport) -> Path:
    suffix = {"pptx": ".pptx", "pdf": ".pdf", "xlsx": ".xlsx"}[report.report_format.value]
    safe_title = "".join(ch if ch.isalnum() else "_" for ch in report.title.lower()).strip("_")
    return get_settings().report_dir / f"{report.id}_{safe_title}{suffix}"


def _project_rows(db: Session, report: GeneratedReport) -> list[Project]:
    query = db.query(Project)
    if "project:" in report.scope:
        project_id = report.scope.split("project:", 1)[1].strip().split()[0]
        project = db.get(Project, project_id)
        return [project] if project else []
    if "account:" in report.scope:
        account_id = report.scope.split("account:", 1)[1].strip().split()[0]
        query = query.filter(Project.account_id == account_id)
    return query.order_by(Project.name).all()


def _scope_value(scope: str, key: str) -> str | None:
    marker = f"{key}:"
    if marker not in scope:
        return None
    return scope.split(marker, 1)[1].strip().split()[0]


def _status_rows(db: Session, report: GeneratedReport, projects: list[Project]) -> list[WeeklyStatus]:
    query = db.query(WeeklyStatus)
    employee_id = _scope_value(report.scope, "employee")
    project_id = _scope_value(report.scope, "project")
    account_id = _scope_value(report.scope, "account")
    period = _scope_value(report.scope, "period")

    if employee_id:
        query = query.filter(WeeklyStatus.employee_id == employee_id)
    if project_id:
        query = query.filter(WeeklyStatus.project_id == project_id)
    elif account_id:
        project_ids = [project.id for project in projects]
        if project_ids:
            query = query.filter(WeeklyStatus.project_id.in_(project_ids))
    rows = query.order_by(WeeklyStatus.week_start.desc()).limit(120).all()
    if period:
        labels = {"daily": "Daily", "weekly": "Weekly", "monthly": "Monthly"}
        expected = labels.get(period.lower())
        if expected:
            rows = [
                row for row in rows
                if str(row.fields.get("reportingFrequency") or row.fields.get("frequency") or "Weekly") == expected
            ]
    return rows[:60]


def _text(value: object | None, fallback: str = "Not reported") -> str:
    if value is None:
        return fallback
    cleaned = str(value).strip()
    return cleaned or fallback


def _short(value: object | None, limit: int = 120, fallback: str = "Not reported") -> str:
    return shorten(_text(value, fallback), width=limit, placeholder="...")


def _clean_llm_text(value: str) -> str:
    cleaned = value.replace("**", "").replace("__", "").replace("###", "").strip()
    return "\n".join(line.strip(" -") for line in cleaned.splitlines() if line.strip())


def _is_meaningful(value: object | None) -> bool:
    if value is None:
        return False
    cleaned = str(value).strip().lower()
    return cleaned not in {"", "none", "no", "no blockers", "n/a", "na", "not applicable"}


def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _set_slide_title(slide, title: str) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = title
        title_shape = slide.shapes.title
    else:
        title_shape = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(9.1), Inches(0.55))
        title_shape.text_frame.text = title
    p = title_shape.text_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = _rgb(BRAND_NAVY)


def _add_textbox(slide, text: str, x: float, y: float, w: float, h: float, size: int = 12, bold: bool = False, color: str = INK_SOFT) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.text = text
    for paragraph in frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = _rgb(color)


def _add_metric_card(slide, label: str, value: str | int, x: float, y: float, w: float = 1.65, h: float = 0.9, accent: str = BRAND_BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(LIGHT_BLUE)
    shape.line.color.rgb = _rgb(BORDER_BLUE)
    tf = shape.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.text = str(value)
    tf.paragraphs[0].font.size = Pt(19)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _rgb(accent)
    p = tf.add_paragraph()
    p.text = label
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = _rgb(INK_SOFT)


def _metrics(statuses: list[WeeklyStatus], projects: list[Project]) -> dict:
    submitted = [s for s in statuses if s.status.value in {"submitted", "approved"}]
    approved = [s for s in statuses if s.status.value == "approved"]
    blockers = [s for s in statuses if _is_meaningful(s.fields.get("blockers")) or str(s.fields.get("overallStatus", "")).lower() == "red"]
    hours = sum(int(s.fields.get("hoursWorked") or 0) for s in statuses)
    completion_values = [int(s.fields.get("completionPercent") or 0) for s in statuses if s.fields.get("completionPercent") is not None]
    avg_completion = round(sum(completion_values) / len(completion_values)) if completion_values else round(sum(p.completion_percent for p in projects) / len(projects)) if projects else 0
    health_counts = {"Green": 0, "Amber": 0, "Red": 0}
    for status in statuses:
        health = str(status.fields.get("overallStatus") or "Green")
        if health in health_counts:
            health_counts[health] += 1
    return {
        "project_count": len(projects),
        "status_count": len(statuses),
        "submitted_count": len(submitted),
        "approved_count": len(approved),
        "blocker_count": len(blockers),
        "hours": hours,
        "avg_completion": avg_completion,
        "health_counts": health_counts,
        "blockers": blockers,
    }


def _fallback_summary(report: GeneratedReport, projects: list[Project], statuses: list[WeeklyStatus]) -> str:
    metrics = _metrics(statuses, projects)
    high_risk_projects = [project.name for project in projects if project.risk.value in {"high", "critical"} or project.health.value == "red"]
    risk_phrase = "; ".join(high_risk_projects[:3]) if high_risk_projects else "no critical project-level risk currently flagged"
    return (
        f"{report.title} covers {metrics['project_count']} project(s) and {metrics['status_count']} status update(s). "
        f"Average completion is {metrics['avg_completion']}%, with {metrics['approved_count']} approved status update(s) and "
        f"{metrics['blocker_count']} blocker or red-health item(s). Current governance attention should focus on {risk_phrase}. "
        "The recommended leadership action is to close blockers, confirm dependencies, and keep weekly approvals current."
    )


def _llm_summary(report: GeneratedReport, projects: list[Project], statuses: list[WeeklyStatus], llm: LLMSelection | None) -> str:
    fallback = _fallback_summary(report, projects, statuses)
    if not llm:
        return fallback
    from app.services.llm import generate_text

    status_digest = "\n".join(
        f"- {_text(status.fields.get('project'), 'Project')} | {status.week_start} | {status.status.value} | "
        f"health={_text(status.fields.get('overallStatus'), 'Green')} | complete={status.fields.get('completionPercent') or 0}% | "
        f"period={_text(status.fields.get('reportingFrequency') or status.fields.get('frequency'), 'Weekly')} | "
        f"update={_short(status.fields.get('achievements'), 140)} | blocker={_short(status.fields.get('blockers'), 100, 'None')}"
        for status in statuses[:16]
    )
    metrics = _metrics(statuses, projects)
    prompt = (
        "Write premium, client-ready PowerPoint copy in 120-150 words for a delivery governance status deck. "
        "Use the exact evidence below; do not invent progress, dates, people, budget, or milestones. "
        "Structure the copy as three concise labeled lines: Delivery position, Risk posture, Leadership actions. "
        "Make the language executive, specific, and analytical. Avoid filler, apologies, markdown, asterisks, and generic phrases such as risks and challenges. "
        "Do not describe contributor status updates as separate projects; if there is one project, say one project with multiple contributors.\n\n"
        f"Report title: {report.title}\nScope: {report.scope}\nRequested period filter: {_scope_value(report.scope, 'period') or 'all'}\n"
        f"Project count: {metrics['project_count']}\nStatus update count: {metrics['status_count']}\n"
        f"Average completion: {metrics['avg_completion']}%\nBlocker/red item count: {metrics['blocker_count']}\n"
        f"Health distribution: {metrics['health_counts']}\n"
        f"Projects: {', '.join(project.name for project in projects) if projects else 'None'}\n"
        f"Contributor status digest:\n{status_digest or 'No status updates available.'}"
    )
    try:
        text, _ = generate_text(llm.provider, prompt, llm.model)
        return _clean_llm_text(_text(text, fallback))
    except Exception:
        return fallback


def generate_report_file(db: Session, report_id: str, llm: LLMSelection | None = None) -> str:
    report = db.get(GeneratedReport, report_id)
    if report is None:
        raise ValueError("Report not found")
    projects = _project_rows(db, report)
    statuses = _status_rows(db, report, projects)
    path = _report_path(report)
    if report.report_format == ReportFormat.PPTX:
        _generate_ppt(path, report, projects, statuses, db, llm)
    elif report.report_format == ReportFormat.PDF:
        _generate_pdf(path, report, projects, statuses, db, llm)
    else:
        _generate_xlsx(path, report, projects, statuses, db)
    report.file_path = str(path)
    report.status = "ready"
    db.commit()
    return str(path)


def _generate_ppt(path: Path, report: GeneratedReport, projects: list[Project], statuses: list[WeeklyStatus], db: Session, llm: LLMSelection | None = None) -> None:
    prs = Presentation(report.template.file_path) if getattr(report, "template", None) and report.template and report.template.file_type == "pptx" else Presentation()

    if len(prs.slides) == 0:
        title = prs.slides.add_slide(prs.slide_layouts[0])
        _set_slide_title(title, report.title)
        if len(title.placeholders) > 1:
            title.placeholders[1].text = "AI-assisted delivery governance report"

    metrics = _metrics(statuses, projects)
    summary_text = _llm_summary(report, projects, statuses, llm)

    summary = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    _set_slide_title(summary, "Executive Delivery Summary")
    _add_textbox(summary, summary_text, 0.65, 1.05, 8.9, 1.45, size=14, color=BRAND_NAVY)
    _add_textbox(summary, f"Scope: {report.scope}", 0.65, 0.75, 8.7, 0.25, size=10, bold=True, color=BRAND_BLUE)
    action_lines = []
    for status in (metrics["blockers"] or statuses[:3])[:3]:
        action_lines.append(f"- {_short(status.fields.get('supportRequired') or status.fields.get('nextWeekPlan'), 120, 'Track next delivery milestone')}")
    _add_textbox(summary, "Leadership Actions\n" + "\n".join(action_lines or ["- Continue weekly governance reviews and approval follow-up."]), 0.75, 2.85, 8.4, 1.35, size=12, color=INK_SOFT)

    dashboard = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    _set_slide_title(dashboard, "Delivery Status Dashboard")
    cards = [
        ("Projects", metrics["project_count"], BRAND_BLUE),
        ("Updates", metrics["status_count"], BRAND_BLUE),
        ("Approved", metrics["approved_count"], SUCCESS),
        ("Blocked", metrics["blocker_count"], DANGER if metrics["blocker_count"] else SUCCESS),
        ("Avg Complete", f"{metrics['avg_completion']}%", WARNING if metrics["avg_completion"] < 70 else SUCCESS),
    ]
    for idx, (label, value, accent) in enumerate(cards):
        _add_metric_card(dashboard, label, value, 0.45 + idx * 1.88, 1.0, accent=accent)

    _add_textbox(dashboard, "Health Distribution", 0.7, 2.25, 2.5, 0.3, size=13, bold=True, color=BRAND_NAVY)
    max_count = max(metrics["health_counts"].values()) if metrics["health_counts"] else 1
    for idx, (health, count) in enumerate(metrics["health_counts"].items()):
        y = 2.75 + idx * 0.48
        _add_textbox(dashboard, health, 0.75, y, 0.85, 0.22, size=9, bold=True)
        width = 0.35 + (count / max(max_count, 1)) * 3.8
        color = {"Green": SUCCESS, "Amber": WARNING, "Red": DANGER}[health]
        bar = dashboard.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.55), Inches(y), Inches(width), Inches(0.22))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(color)
        bar.line.color.rgb = _rgb(color)
        _add_textbox(dashboard, str(count), 5.5, y, 0.45, 0.22, size=9, bold=True)
    _add_textbox(dashboard, "Interpretation\nGreen indicates on-track execution. Amber requires management follow-up. Red or blocker-backed items should be treated as escalation candidates for the next governance review.", 6.25, 2.35, 3.0, 1.6, size=11)

    table_slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    _set_slide_title(table_slide, "Team Status Detail")
    rows = max(min(len(statuses), 8) + 1, 2)
    table = table_slide.shapes.add_table(rows, 6, Inches(0.25), Inches(1.05), Inches(9.5), Inches(4.9)).table
    headers = ["Person", "Project", "Cycle", "Health", "Complete", "Key Update"]
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(BRAND_NAVY)
        cell.text_frame.paragraphs[0].font.color.rgb = _rgb("#FFFFFF")
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(9)
    for row_idx, status in enumerate(statuses[:8], start=1):
        employee = db.get(Employee, status.employee_id)
        project = db.get(Project, status.project_id) if status.project_id else None
        values = [
            employee.name if employee else "-",
            project.name if project else _short(status.fields.get("project"), 24, "-"),
            f"{_text(status.fields.get('reportingFrequency') or status.fields.get('frequency'), 'Weekly')} / {status.week_start.strftime('%d %b %Y')}",
            _text(status.fields.get("overallStatus"), status.status.value),
            f"{status.fields.get('completionPercent') or 0}%",
            _short(status.fields.get("achievements"), 90, "No narrative submitted"),
        ]
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = _rgb(BRAND_NAVY if col_idx == 0 else INK_SOFT)

    risk_slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    _set_slide_title(risk_slide, "Risks, Decisions, Next Actions")
    risk_items = metrics["blockers"] or statuses[:5]
    y = 1.05
    for status in risk_items[:5]:
        project_name = _text(status.fields.get("project"), "Project")
        risk = _short(status.fields.get("risks") or status.fields.get("blockers"), 120, "No major risk reported")
        action = _short(status.fields.get("supportRequired") or status.fields.get("nextWeekPlan"), 120, "Continue tracking next milestone")
        _add_textbox(risk_slide, project_name, 0.65, y, 2.1, 0.25, size=10, bold=True, color=BRAND_BLUE)
        _add_textbox(risk_slide, f"Risk: {risk}\nAction: {action}", 2.75, y, 6.7, 0.55, size=10, color=INK_SOFT)
        y += 0.75
    if not risk_items:
        _add_textbox(risk_slide, "No submitted risks or blockers are available for this scope.", 0.65, 1.2, 8.7, 0.45, size=13)

    prs.save(path)


def _generate_pdf(path: Path, report: GeneratedReport, projects: list[Project], statuses: list[WeeklyStatus], db: Session, llm: LLMSelection | None = None) -> None:
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(path), pagesize=A4, rightMargin=28, leftMargin=28, topMargin=28, bottomMargin=28)
    metrics = _metrics(statuses, projects)
    summary_text = _llm_summary(report, projects, statuses, llm)
    story = [Paragraph(report.title, styles["Title"]), Paragraph("AI-assisted delivery governance report", styles["Normal"]), Spacer(1, 12), Paragraph(summary_text, styles["BodyText"]), Spacer(1, 14)]

    kpi_data = [["Projects", "Updates", "Approved", "Blocked", "Avg Complete"], [metrics["project_count"], metrics["status_count"], metrics["approved_count"], metrics["blocker_count"], f"{metrics['avg_completion']}%"]]
    kpi_table = Table(kpi_data, repeatRows=1)
    kpi_table.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1D4ED8")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white), ("ALIGN", (0, 0), (-1, -1), "CENTER"), ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#BFDBFE")), ("FONTNAME", (0, 1), (-1, 1), "Helvetica-Bold")]))
    story.extend([kpi_table, Spacer(1, 16)])

    data = [["Project", "Account", "Phase", "Health", "Risk", "Completion"]]
    for project in projects:
        account = db.get(Account, project.account_id)
        data.append([project.name, account.name if account else "-", project.phase.value, project.health.value, project.risk.value, f"{project.completion_percent}%"])
    table = Table(data, repeatRows=1)
    table.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0F172A")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white), ("GRID", (0, 0), (-1, -1), 0.25, colors.grey), ("VALIGN", (0, 0), (-1, -1), "TOP")]))
    story.append(table)

    if statuses:
        story.extend([Spacer(1, 18), Paragraph("Status Details", styles["Heading2"])])
        status_data = [["Person", "Cycle", "Status", "Health", "Key Update"]]
        for status in statuses[:14]:
            employee = db.get(Employee, status.employee_id)
            status_data.append([employee.name if employee else "-", f"{_text(status.fields.get('reportingFrequency') or status.fields.get('frequency'), 'Weekly')} / {status.week_start.strftime('%d %b %Y')}", status.status.value, _text(status.fields.get("overallStatus"), "-"), _short(status.fields.get("achievements"), 95, "No update")])
        status_table = Table(status_data, repeatRows=1)
        status_table.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2563EB")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white), ("GRID", (0, 0), (-1, -1), 0.25, colors.grey), ("VALIGN", (0, 0), (-1, -1), "TOP")]))
        story.append(status_table)
    doc.build(story)

    if getattr(report, "template", None) and report.template and report.template.file_type == "pdf":
        try:
            from pypdf import PdfReader, PdfWriter
        except ImportError as exc:
            raise RuntimeError("PDF template support requires pypdf to be installed") from exc
        final_writer = PdfWriter()
        template_reader = PdfReader(report.template.file_path)
        output_reader = PdfReader(str(path))
        for page in template_reader.pages:
            final_writer.add_page(page)
        for page in output_reader.pages:
            final_writer.add_page(page)
        with path.open("wb") as output_file:
            final_writer.write(output_file)


def _generate_xlsx(path: Path, report: GeneratedReport, projects: list[Project], statuses: list[WeeklyStatus], db: Session) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Delivery Governance"
    ws.append(["Project", "Account", "Phase", "Health", "Risk", "Budget Used", "Budget Total", "Completion"])
    for project in projects:
        account = db.get(Account, project.account_id)
        ws.append([project.name, account.name if account else "-", project.phase.value, project.health.value, project.risk.value, float(project.budget_used), float(project.budget_total), project.completion_percent])
    for col in ws.columns:
        ws.column_dimensions[col[0].column_letter].width = 18
    status_ws = wb.create_sheet("Status Details")
    status_ws.append(["Employee", "Project", "Reporting Cycle", "Period Start", "Submission Status", "Health", "Completion", "Hours", "Achievements", "Blockers", "Risks", "Next Plan"])
    for status in statuses:
        employee = db.get(Employee, status.employee_id)
        project = db.get(Project, status.project_id) if status.project_id else None
        status_ws.append([employee.name if employee else "-", project.name if project else _text(status.fields.get("project"), "-"), _text(status.fields.get("reportingFrequency") or status.fields.get("frequency"), "Weekly"), status.week_start.isoformat(), status.status.value, _text(status.fields.get("overallStatus"), "-"), status.fields.get("completionPercent") or 0, status.fields.get("hoursWorked") or 0, _text(status.fields.get("achievements"), ""), _text(status.fields.get("blockers"), ""), _text(status.fields.get("risks"), ""), _text(status.fields.get("nextWeekPlan"), "")])
    for col in status_ws.columns:
        status_ws.column_dimensions[col[0].column_letter].width = 22
    wb.save(path)



