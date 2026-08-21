from pathlib import Path
from types import SimpleNamespace

from pptx import Presentation
from pptx.util import Inches

from app.ai.ppt_mapping import GeminiMappingConfigurationError, GeminiMappingError, PPTMapping, map_with_gemini
from app.ai.gemini_response import GeminiResponseError, extract_gemini_text, parse_gemini_json
from app.ai.template_analysis import analyze_template
from app.core.config import get_settings
from app.reports.generator import _populate_template


def _template(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
    title.text = "{{PROJECT_NAME}}"
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(2))
    body.text = "{{ACHIEVEMENTS}}"
    prs.save(path)


def test_template_analysis_extracts_structure(tmp_path):
    path = tmp_path / "account-template.pptx"
    _template(path)
    structure = analyze_template(path)
    assert structure.slide_count == 1
    assert structure.slides[0].elements[0].text == "{{PROJECT_NAME}}"
    assert structure.slide_width > 0


def test_gemini_mapping_validates_structured_response(monkeypatch, tmp_path):
    path = tmp_path / "account-template.pptx"
    _template(path)

    class FakeModels:
        def generate_content(self, **_kwargs):
            return SimpleNamespace(text='{"slides":[{"slide_index":0,"element_fields":{"slide_0_shape_0":["PROJECT_NAME"],"slide_0_shape_1":["ACHIEVEMENTS"]}}]}')

    monkeypatch.setattr("google.genai.Client", lambda **_kwargs: SimpleNamespace(models=FakeModels()))
    monkeypatch.setattr("app.ai.ppt_mapping.get_settings", lambda: SimpleNamespace(gemini_api_key="configured", gemini_default_model="gemini-2.5-flash"))
    mapping = map_with_gemini(analyze_template(path), {"project": [{"name": "Project A1"}]})
    assert isinstance(mapping, PPTMapping)
    assert mapping.slides[0].element_fields["slide_0_shape_1"] == ["ACHIEVEMENTS"]


def test_gemini_text_extractor_ignores_thought_signature():
    response = SimpleNamespace(
        candidates=[SimpleNamespace(content=SimpleNamespace(parts=[
            SimpleNamespace(text='{"account_name":"Account A",'),
            SimpleNamespace(thought_signature=b"opaque-model-metadata"),
            SimpleNamespace(text='"project_name":"Project A1"}'),
        ]))],
    )
    assert extract_gemini_text(response) == '{"account_name":"Account A","project_name":"Project A1"}'
    assert parse_gemini_json(response)["project_name"] == "Project A1"


def test_gemini_mapping_accepts_fenced_json_with_non_text_parts(monkeypatch, tmp_path):
    path = tmp_path / "account-template.pptx"
    _template(path)

    class FakeModels:
        def generate_content(self, **_kwargs):
            return SimpleNamespace(candidates=[SimpleNamespace(content=SimpleNamespace(parts=[
                SimpleNamespace(text='```json\n{"slides":[{"slide_index":0,"element_fields":{"slide_0_shape_0":["PROJECT_NAME"]}}]}\n```'),
                SimpleNamespace(thought_signature="ignored"),
            ]))])

    monkeypatch.setattr("google.genai.Client", lambda **_kwargs: SimpleNamespace(models=FakeModels()))
    monkeypatch.setattr("app.ai.ppt_mapping.get_settings", lambda: SimpleNamespace(gemini_api_key="configured", gemini_default_model="gemini-3.5-flash"))
    mapping = map_with_gemini(analyze_template(path), {})
    assert mapping.slides[0].element_fields["slide_0_shape_0"] == ["PROJECT_NAME"]


def test_malformed_gemini_mapping_response_is_actionable(monkeypatch, tmp_path):
    path = tmp_path / "account-template.pptx"
    _template(path)

    class FakeModels:
        def generate_content(self, **_kwargs):
            return SimpleNamespace(candidates=[SimpleNamespace(content=SimpleNamespace(parts=[SimpleNamespace(text="not json")]))])

    monkeypatch.setattr("google.genai.Client", lambda **_kwargs: SimpleNamespace(models=FakeModels()))
    monkeypatch.setattr("app.ai.ppt_mapping.get_settings", lambda: SimpleNamespace(gemini_api_key="configured", gemini_default_model="gemini-3.5-flash"))
    try:
        map_with_gemini(analyze_template(path), {})
    except RuntimeError as exc:
        assert "invalid PPT mapping response" in str(exc)
    else:
        raise AssertionError("Expected an actionable malformed JSON error")


def test_gemini_response_without_text_parts_is_rejected():
    response = SimpleNamespace(candidates=[SimpleNamespace(content=SimpleNamespace(parts=[SimpleNamespace(thought_signature="ignored")]))])
    try:
        extract_gemini_text(response)
    except GeminiResponseError as exc:
        assert "no text parts" in str(exc)
    else:
        raise AssertionError("Expected missing text parts to be rejected")


def test_gemini_mapping_without_text_parts_is_actionable(monkeypatch, tmp_path):
    path = tmp_path / "account-template.pptx"
    _template(path)

    class FakeModels:
        def generate_content(self, **_kwargs):
            return SimpleNamespace(candidates=[SimpleNamespace(content=SimpleNamespace(parts=[SimpleNamespace(thought_signature="ignored")]))])

    monkeypatch.setattr("google.genai.Client", lambda **_kwargs: SimpleNamespace(models=FakeModels()))
    monkeypatch.setattr("app.ai.ppt_mapping.get_settings", lambda: SimpleNamespace(gemini_api_key="configured", gemini_default_model="gemini-3.5-flash"))
    try:
        map_with_gemini(analyze_template(path), {})
    except RuntimeError as exc:
        assert "invalid PPT mapping response" in str(exc)
    else:
        raise AssertionError("Expected missing Gemini text to fail clearly")


def test_shape_targeted_mapping_populates_all_status_fields(tmp_path):
    path = tmp_path / "labelled-account-template.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fields = [
        "ACCOUNT_NAME", "PROJECT_NAME", "OVERALL_STATUS", "COMPLETION",
        "HOURS", "ACHIEVEMENTS", "BLOCKERS", "RISKS", "NEXT_WEEK_PLAN",
    ]
    for index, field in enumerate(fields):
        box = slide.shapes.add_textbox(Inches(1), Inches(0.4 + index * 0.5), Inches(8), Inches(0.35))
        box.text = field
    prs.save(path)

    structure = analyze_template(path)
    element_fields = {
        element.id: [fields[index]]
        for index, element in enumerate(structure.slides[0].elements)
    }
    mapping = PPTMapping(slides=[{"slide_index": 0, "element_fields": element_fields}])
    values = {
        "ACCOUNT_NAME": "Arbitrary Account",
        "PROJECT_NAME": "Arbitrary Project",
        "OVERALL_STATUS": "Amber",
        "COMPLETION": "73%",
        "HOURS": "41",
        "ACHIEVEMENTS": "Completed integration testing",
        "BLOCKERS": "Waiting for upstream approval",
        "RISKS": "Approval delay may affect release",
        "NEXT_WEEK_PLAN": "Start user acceptance testing",
    }

    populated = Presentation(str(path))
    assert _populate_template(populated, values, mapping)
    output = tmp_path / "populated.pptx"
    populated.save(output)

    reopened = Presentation(str(output))
    text = "\n".join(shape.text for slide in reopened.slides for shape in slide.shapes if getattr(shape, "text", ""))
    for value in values.values():
        assert value in text


def test_token_based_mapping_remains_supported(tmp_path):
    path = tmp_path / "token-template.pptx"
    _template(path)
    prs = Presentation(str(path))

    assert _populate_template(prs, {
        "PROJECT_NAME": "Token Project",
        "ACHIEVEMENTS": "Token achievement",
    })
    text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if getattr(shape, "text", ""))
    assert "Token Project" in text
    assert "Token achievement" in text


def test_gemini_mapping_requires_api_key(monkeypatch):
    monkeypatch.setattr("app.ai.ppt_mapping.get_settings", lambda: SimpleNamespace(gemini_api_key=None, gemini_default_model="gemini-3.5-flash"))
    try:
        map_with_gemini(SimpleNamespace(model_dump_json=lambda: "{}"), {})
    except GeminiMappingConfigurationError as exc:
        assert "not configured" in str(exc)
    else:
        raise AssertionError("Expected missing Gemini configuration to fail")


def test_invalid_gemini_model_is_controlled(monkeypatch):
    class FakeModels:
        def generate_content(self, **_kwargs):
            raise ValueError("model does not exist")

    monkeypatch.setattr("google.genai.Client", lambda **_kwargs: SimpleNamespace(models=FakeModels()))
    monkeypatch.setattr("app.ai.ppt_mapping.get_settings", lambda: SimpleNamespace(gemini_api_key="configured", gemini_default_model="invalid-model"))
    try:
        map_with_gemini(SimpleNamespace(model_dump_json=lambda: "{}"), {})
    except GeminiMappingError as exc:
        assert str(exc) == "Gemini PPT mapping failed."
        assert "configured" not in str(exc)
    else:
        raise AssertionError("Expected invalid Gemini model to fail safely")


def test_empty_gemini_mapping_is_rejected(monkeypatch, tmp_path):
    path = tmp_path / "account-template.pptx"
    _template(path)

    class FakeModels:
        def generate_content(self, **_kwargs):
            return SimpleNamespace(text='{"slides":[{"slide_index":0,"element_fields":{}}]}')

    monkeypatch.setattr("google.genai.Client", lambda **_kwargs: SimpleNamespace(models=FakeModels()))
    monkeypatch.setattr("app.ai.ppt_mapping.get_settings", lambda: SimpleNamespace(gemini_api_key="configured", gemini_default_model="gemini-3.5-flash"))
    try:
        map_with_gemini(analyze_template(path), {})
    except RuntimeError as exc:
        assert "empty PPT mapping" in str(exc)
    else:
        raise AssertionError("Expected empty mapping to be rejected")
